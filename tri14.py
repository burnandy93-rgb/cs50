from pypdf import PdfReader
from pptx import Presentation
import os

pdf_file = r"C:\Users\USER\IdeaProjects\untitled1trial\command.pdf"
ppt_file = r"C:\Users\USER\IdeaProjects\untitled1trial\output.pptx"

if os.path.exists(pdf_file):
    reader = PdfReader(pdf_file)
    ppt = Presentation()
    for page in reader.pages:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = "PDF Content"
        slide.placeholders[1].text = page.extract_text()[:500]
    ppt.save(ppt_file)
    print("PDF converted to PowerPoint successfully")
else:
    print(f"File not found: {pdf_file}")
